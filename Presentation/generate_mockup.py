#!/usr/bin/env python3
"""Макет презентации на базе шаблона МАИ (Теленков КД ВКР.pptx)."""

from __future__ import annotations

import copy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

ROOT = Path(__file__).resolve().parent.parent
PRESENTATION_DIR = Path(__file__).resolve().parent
TEMPLATE = PRESENTATION_DIR / "Пример" / "Теленков КД ВКР.pptx"
OUT = PRESENTATION_DIR / "DCT_Защита_макет.pptx"
OUT_FALLBACK = PRESENTATION_DIR / "DCT_Защита_макет_новый.pptx"

# Фирменный цвет (из макета кафедры)
C_PRIMARY = RGBColor(0x00, 0x95, 0xDA)  # #0095DA
C_TEXT = RGBColor(0x22, 0x22, 0x22)
C_MUTED = RGBColor(0x66, 0x66, 0x66)

THESIS_TITLE = (
    "Разработка программно-аппаратного комплекса отслеживания "
    "малоразмерных FPV-БВС по RC-каналам и видеопотоку курсовой камеры (DCT)"
)

TITLE_AUTHOR_PLACEHOLDER = (
    "Выполнил:\n"
    "студент группы [М3О-XXX-XX] [Фамилия И. О.]\n"
    "Научный руководитель:\n"
    "[учёная степень, звание, ФИО]\n"
    "Консультант:\n"
    "[при наличии]"
)

THANKS_CONTACT_PLACEHOLDER = (
    "[Фамилия Имя Отчество]\n"
    "[группа М3О-XXX-XX]\n"
    "[email@example.com]"
)

# Позиции для правой колонки (визуал), EMU
IMG_LEFT = 6_500_000
IMG_TOP = 1_500_000
IMG_WIDTH = 5_500_000
IMG_HEIGHT = 4_800_000


def delete_slide(prs: Presentation, index: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)
    slide_id = slides[index]
    r_id = slide_id.rId
    prs.part.drop_rel(r_id)
    slide_id_list.remove(slide_id)


def clone_shape_elements(source_slide) -> list:
    return [copy.deepcopy(sh.element) for sh in source_slide.shapes]


def add_slide_from_elements(prs: Presentation, elements: list):
    layout = prs.slide_layouts[6]  # Пустой слайд
    slide = prs.slides.add_slide(layout)
    for el in elements:
        slide.shapes._spTree.insert_element_before(copy.deepcopy(el), "p:extLst")
    return slide


def find_shape(slide, *needles: str):
    for sh in slide.shapes:
        name = sh.name or ""
        for n in needles:
            if n in name:
                return sh
    return None


def set_shape_text(shape, text: str, font_size: int | None = None, color: RGBColor | None = None) -> None:
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    if font_size:
        p.font.size = Pt(font_size)
    if color:
        p.font.color.rgb = color
    p.font.name = "Arial"


def set_title_style(shape) -> None:
    if not shape or not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for run in p.runs:
            run.font.color.rgb = C_PRIMARY
            run.font.bold = True


def set_body_bullets(shape, bullets: list[str]) -> None:
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    if not bullets:
        tf.text = ""
        return
    tf.text = bullets[0]
    for line in bullets[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)
        p.font.name = "Arial"
        p.font.color.rgb = C_TEXT
        p.space_after = Pt(6)
    for p in tf.paragraphs:
        p.font.size = Pt(18)
        p.font.name = "Arial"
        p.font.color.rgb = C_TEXT


def add_image_or_placeholder(slide, hint: str, image_path: Path | None) -> None:
    if image_path and image_path.is_file():
        try:
            slide.shapes.add_picture(
                str(image_path),
                IMG_LEFT,
                IMG_TOP,
                width=IMG_WIDTH,
                height=IMG_HEIGHT,
            )
            return
        except Exception:
            pass
    box = slide.shapes.add_textbox(
        IMG_LEFT,
        IMG_TOP + IMG_HEIGHT - 400_000,
        IMG_WIDTH,
        400_000,
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f"[Визуал: {hint[:90]}]"
    p.font.size = Pt(10)
    p.font.color.rgb = C_MUTED
    p.alignment = PP_ALIGN.CENTER


def set_slide_number(shape, number: int) -> None:
    if shape and shape.has_text_frame:
        set_shape_text(shape, str(number), font_size=14, color=C_MUTED)


def speaker_notes(slide, essence: str, timing: str, block: str | None = None) -> None:
    parts = [f"Суть: {essence}", f"Время: {timing}"]
    if block:
        parts.append(f"Блок: {block}")
    slide.notes_slide.notes_text_frame.text = "\n".join(parts)


def update_title_slide(slide) -> None:
    sh_title = find_shape(slide, "Заголовок 3")
    sh_author = find_shape(slide, "Подзаголовок 2")
    if sh_title:
        set_shape_text(sh_title, THESIS_TITLE, font_size=28)
        set_title_style(sh_title)
    if sh_author:
        set_shape_text(sh_author, TITLE_AUTHOR_PLACEHOLDER, font_size=16)


def update_thanks_slide(slide, slide_num: int) -> None:
    sh_title = find_shape(slide, "Заголовок 4", "Заголовок")
    sh_text = find_shape(slide, "Текст 5", "Текст")
    sh_num = find_shape(slide, "Номер слайда")
    if sh_title:
        set_title_style(sh_title)
    if sh_text:
        set_shape_text(sh_text, THANKS_CONTACT_PLACEHOLDER, font_size=20)
    set_slide_number(sh_num, slide_num)


def fill_content_slide(
    slide,
    data: dict,
    presentation_slide_num: int,
    narrow_body: bool = True,
) -> None:
    sh_title = find_shape(slide, "Заголовок 8", "Заголовок")
    sh_body = find_shape(slide, "TextBox 4", "Объект 9", "Объект")
    sh_num = find_shape(slide, "Номер слайда")
    sh_obj_short = find_shape(slide, "Объект 9")

    if sh_title:
        set_shape_text(sh_title, data["title"], font_size=28)
        set_title_style(sh_title)

    bullets = data["bullets"]
    if sh_obj_short and sh_obj_short != sh_body:
        set_shape_text(sh_obj_short, "")
    if sh_body:
        set_body_bullets(sh_body, bullets)
        if narrow_body and (data.get("image") or data.get("visual_hint")):
            # Уже текст слева за счёт ширины TextBox — картинка справа перекрывает часть
            pass

    set_slide_number(sh_num, presentation_slide_num)

    img_rel = data.get("image")
    img_path = ROOT / img_rel if img_rel else None
    add_image_or_placeholder(slide, data.get("visual_hint", ""), img_path)

    speaker_notes(
        slide,
        data.get("essence", ""),
        data.get("timing", ""),
        data.get("block"),
    )


SLIDES = [
    {
        "num": 1,
        "block": "Блок 1–2 · ввод",
        "title": "Проблема FPV-мониторинга",
        "timing": "0:30",
        "essence": "Полёт без NAV возможен; координаты для мониторинга нужны; борт не дорабатываем.",
        "bullets": [
            "Малоразмерный FPV-БВС может летать без бортовой навигации",
            "Без координат нет контроля трассы и передачи в системы мониторинга",
            "Уже доступны: RC-каналы пилота и видео с курсовой камеры",
            "Задача: оценить положение на наземной станции, не дооснащая дрон",
        ],
        "visual_hint": "Схема: Пилот → БВС → ПК (DCT) → мониторинг",
    },
    {
        "num": 2,
        "block": "Блок 1–2 · ввод",
        "title": "Почему не GPS / UWB / SLAM",
        "timing": "0:30",
        "essence": "Типовые методы требуют борта или инфраструктуры.",
        "bullets": [
            "GPS/ГЛОНАСС — приёмник и телеметрия на борту",
            "UWB, motion capture — маяки и подготовка площадки",
            "Внешние камеры — поле зрения, освещение, скорость",
            "SLAM / визуальная одометрия — нагрузка и нестабильная FPV-сцена",
        ],
        "visual_hint": "Табл. 1 из ВКР (4–5 строк)",
    },
    {
        "num": 3,
        "block": "Блок 1–2 · ввод",
        "title": "Решение: наземный комплекс DCT",
        "timing": "0:30",
        "essence": "RC + видео на ПК; локализация по трассе; Record/Replay.",
        "bullets": [
            "Сбор и синхронизация RC + видео на ПК",
            "Локализация относительно заданной трассы (референс + PF)",
            "Режимы Record (запись) и Replay (анализ без нового полёта)",
            "Выход: координаты во внешнюю систему мониторинга",
        ],
        "visual_hint": "Рис. 2 — C4 System Context",
    },
    {
        "num": 4,
        "block": "Блок 1–2 · ввод",
        "title": "Цель и задачи проекта",
        "timing": "0:30",
        "essence": "5 задач; критерий p90 < 15 м.",
        "bullets": [
            "Цель: комплекс сбора, локализации и передачи координат БВС",
            "(1) оборудование и сбор данных",
            "(2) архитектура (3) данные sim/real",
            "(4) отслеживание (5) интеграция с мониторингом",
            "Критерий исследования: p90 ошибки < 15 м",
        ],
        "visual_hint": "Табл. 3 из ВКР",
    },
    {
        "num": 5,
        "block": "Блок 3 · теория",
        "title": "Архитектура комплекса",
        "timing": "0:35",
        "essence": "4 блока; вычисления на земле.",
        "bullets": [
            "Четыре блока: RC, видео, DCT, интеграция с мониторингом",
            "Внутри DCT: запись → воспроизведение → алгоритмы",
            "Вычисления на наземной станции, борт без NAV-модуля",
        ],
        "visual_hint": "Рис. 3 — C4 Container",
    },
    {
        "num": 6,
        "block": "Блок 3 · теория",
        "title": "Аппаратура и синхронизация",
        "timing": "0:30",
        "essence": "Типовой FPV-стенд; единая временная шкала.",
        "bullets": [
            "RC: EdgeTX/ELRS → приёмник → ESP32 → COM → DCT",
            "Видео: FPV-камера → HDZero → захват на ПК",
            "Все потоки на единой временной шкале (wall-clock + кадры)",
        ],
        "visual_hint": "Рис. 7–8: rc_chain + video_chain",
    },
    {
        "num": 7,
        "block": "Блок 3 · теория",
        "title": "Модель трассы и состояния",
        "timing": "0:30",
        "essence": "Референс .npz; состояние (s, v).",
        "bullets": [
            "Трасса: ворота + референсная траектория (файл .npz)",
            "Состояние фильтра: (s, v) — положение вдоль дуги (м) и скорость (м/с)",
            "RC-признаки сопоставляются с узлами референса (режим RC+Rate)",
        ],
        "visual_hint": "Рис. 11 — ворота и референс",
    },
    {
        "num": 8,
        "block": "Блок 3 · теория",
        "title": "Particle Filter",
        "timing": "0:30",
        "essence": "PF — основной алгоритм; KFLayer2 поверх.",
        "bullets": [
            "Частицы (s, v): predict → update по RC → resample",
            "Циклический круг: положение по модулю длины трассы",
            "Поверх PF — сглаживание KFLayer2 (не замена PF)",
        ],
        "visual_hint": "Рис. 45 / pf_cycle.drawio",
    },
    {
        "num": 9,
        "block": "Блок 3 · теория",
        "title": "Камера в контуре локализации",
        "timing": "0:25",
        "essence": "YOLO→PnP→inject; модуль опционален.",
        "bullets": [
            "YOLO (ворота) → PnP → CameraObservation",
            "inject_position_observation — байесовское обновление весов PF",
            "Основной контур работает без камеры",
        ],
        "visual_hint": "Рис. 46 / camera_observation.drawio",
    },
    {
        "num": 10,
        "block": "Блок 4 · реализация",
        "title": "Структура программного комплекса DCT",
        "timing": "0:40",
        "essence": "Модули; воспроизводимая сессия.",
        "bullets": [
            "Модули: session, replay, profiles, localization, camera, MAVLink",
            "Сессия — самодостаточная папка (RC, видео, meta, observations)",
            "Сбор данных отделён от алгоритмических экспериментов",
        ],
        "visual_hint": "Рис. 49 — C4 Component / Прил. А",
    },
    {
        "num": 11,
        "block": "Блок 4 · реализация",
        "title": "Интерфейс пользователя",
        "timing": "0:40",
        "essence": "Record и Replay с HUD и картой.",
        "bullets": [
            "Record: пилот/дрон/трасса, HUD, карта, видео, запись",
            "Replay: повторный анализ и отладка локализации",
            "Эксперименты без повторного полёта",
        ],
        "visual_hint": "Рис. 9–10 — скриншоты GUI",
    },
    {
        "num": 12,
        "block": "Блок 4 · реализация",
        "title": "Интеграция с «Небосвод»",
        "timing": "0:40",
        "essence": "Геопривязка; MAVLink; p90 vs p95/HEPU.",
        "bullets": [
            "Локальные координаты → WGS84 (origin, x, z)",
            "MAVLink UDP → Nebosvod Connect → «Небосвод»",
            "p90 — для алгоритмов; p95 — для HEPU",
        ],
        "visual_hint": "Рис. 40 — integration_monitoring.drawio",
    },
    {
        "num": 13,
        "block": "Блок 5 · эксперименты",
        "title": "Методика экспериментов",
        "timing": "0:35",
        "essence": "p90; 12 экспериментов в 4 блоках.",
        "bullets": [
            "Метрика: p90_err_m (90-й перцентиль ошибки, м)",
            "Критерий: p90 < 15 м",
            "Блоки: метод (1–4) → устойчивость (5–9) → real (10) → камера (11–12)",
        ],
        "visual_hint": "Табл. 13 — воронка 4 блоков",
    },
    {
        "num": 14,
        "block": "Блок 5 · эксперименты",
        "title": "Выбор и настройка метода (эксп. 1–4)",
        "timing": "0:45",
        "essence": "PF; cross-drone p90 = 7,8 м.",
        "bullets": [
            "Эксп. 1: Particle Filter лучше HMM, DTW, NN",
            "Эксп. 2–3: режим RC+Rate, веса каналов",
            "Эксп. 4: obs_sigma=2, pnv=8 → cross-drone p90 = 7,8 м",
        ],
        "visual_hint": "Рис. 14, 18–19",
        "image": "tools/exp0_benchmark/plots/method_comparison.png",
    },
    {
        "num": 15,
        "block": "Блок 5 · эксперименты",
        "title": "Устойчивость и переносимость (эксп. 5–9)",
        "timing": "0:45",
        "essence": "Качество референса важнее same-drone.",
        "bullets": [
            "Важнее качество референса, чем совпадение дрона",
            "Кросс-пилот: переносимость с ограничениями",
            "Новые трассы: mini-sweep; дрейф по времени не выявлен",
        ],
        "visual_hint": "Рис. 20, 24",
        "image": "tools/exp4_reference/plots/same_vs_cross_drone_ref.png",
    },
    {
        "num": 16,
        "block": "Блок 5 · эксперименты",
        "title": "Реальный трек и камера (эксп. 10–12)",
        "timing": "0:50",
        "essence": "Real track OK; камера offline.",
        "bullets": [
            "Эксп. 10: RC/PF вне симулятора (real track)",
            "Эксп. 11: моделирование слияния — снижение p90",
            "Эксп. 12: YOLO/PnP + CamKF; ограничение: offline",
        ],
        "visual_hint": "Рис. 30–31, 36",
        "image": "tools/exp9_realworld/plots/variant_comparison.png",
    },
    {
        "num": 17,
        "block": "Блок 5 · эксперименты",
        "title": "Итоги экспериментальной части",
        "timing": "0:35",
        "essence": "Критерий выполнен; базовая конфигурация зафиксирована.",
        "bullets": [
            "p90 < 15 м — выполнен в ключевых сценариях",
            "База: RC+Rate, PF, pnv=8, снижение throttle",
            "Камера — усиление; нужны real-time и больший датасет",
        ],
        "visual_hint": "3 плашки с цифрами",
        "image": "tools/exp3_hyperparam/plots/optimal_comparison.png",
    },
    {
        "num": 18,
        "block": "Блок 6 · заключение",
        "title": "Заключение",
        "timing": "1:00",
        "essence": "Итог, ограничения, публикации.",
        "bullets": [
            "Создан DCT: сбор, синхронизация, PF, камера, MAVLink",
            "Доказано: RC/PF; p90 < 15 м; реальный трек",
            "Ограничения: геопривязка, MAVLink, камера offline, p95/HEPU",
            "Публикации / выступления / акты: [вставить факты]",
        ],
        "visual_hint": "Без графиков",
    },
]


def main() -> None:
    if not TEMPLATE.is_file():
        raise FileNotFoundError(f"Шаблон не найден: {TEMPLATE}")

    prs = Presentation(str(TEMPLATE))

    # Прототипы до удаления: слайд 3 (контент), слайд 17 (спасибо)
    content_proto = clone_shape_elements(prs.slides[2])
    thanks_proto = clone_shape_elements(prs.slides[16])

    # Оставить только титульный (индекс 0)
    for i in range(len(prs.slides) - 1, 0, -1):
        delete_slide(prs, i)

    update_title_slide(prs.slides[0])

    # 18 контентных слайдов (нумерация 2–19)
    for i, data in enumerate(SLIDES):
        slide = add_slide_from_elements(prs, content_proto)
        fill_content_slide(slide, data, presentation_slide_num=i + 2)

    # Финальный слайд (20)
    thanks_slide = add_slide_from_elements(prs, thanks_proto)
    update_thanks_slide(thanks_slide, slide_num=20)

    out_path = OUT
    try:
        prs.save(out_path)
    except PermissionError:
        out_path = OUT_FALLBACK
        prs.save(out_path)
        print("NOTE: основной файл занят (закройте PowerPoint). Сохранено в:", out_path.name)
    print(f"Saved: {out_path}")
    print(f"Slides: {len(prs.slides)} (title + {len(SLIDES)} content + thanks)")
    print(f"Template: {TEMPLATE.name}")
    print(f"Color: #0095DA")


if __name__ == "__main__":
    main()
